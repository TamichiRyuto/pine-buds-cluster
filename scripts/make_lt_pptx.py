#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build slides/lt-ear-parallel.pptx from the deck defined in slides/lt-ear-parallel.md.

Theme-less on purpose: blank layout, white ground, black text, grey rules.
Structure comes from shapes and text boxes, not from a PowerPoint design theme.
Sized for projection: body text is 13 pt and up, tables 13.5 pt, code 13 pt.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "slides", "lt-ear-parallel.pptx")

# Drop a real photo at slides/assets/pinebuds.<ext> and it is embedded on the hardware slide.
PHOTO = None
for _ext in ("jpg", "jpeg", "png", "JPG", "JPEG", "PNG"):
    _cand = os.path.join(ROOT, "slides", "assets", "pinebuds." + _ext)
    if os.path.exists(_cand):
        PHOTO = _cand
        break

W, H = 13.333, 7.5
M = 0.55
CW = W - 2 * M
JP = "Meiryo"
MONO = "Consolas"

INK = RGBColor(0x11, 0x11, 0x11)
MUTED = RGBColor(0x4A, 0x4A, 0x4A)
FAINT = RGBColor(0x80, 0x80, 0x80)
RULE = RGBColor(0xBF, 0xBF, 0xBF)
BOX = RGBColor(0x9E, 0x9E, 0x9E)
FILL = RGBColor(0xF4, 0xF4, 0xF4)
FILL2 = RGBColor(0xE6, 0xE6, 0xE6)
CODEBG = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OK = RGBColor(0x1E, 0x6B, 0x3A)
NG = RGBColor(0xA3, 0x2B, 0x22)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- primitives
def font(run, name=JP, size=15, bold=False, color=INK, italic=False):
    f = run.font
    f.name, f.size, f.bold, f.italic = name, Pt(size), bold, italic
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def text(slide, x, y, w, h, lines, size=15, name=JP, color=INK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=5, pad=0.07, lh=None):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(pad)
    tf.margin_top = tf.margin_bottom = Inches(pad * 0.6)
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, item in enumerate(lines):
        # A list item is one paragraph made of several differently styled runs.
        runs = item if isinstance(item, list) else [item]
        pov = runs[0][1] if isinstance(runs[0], tuple) else {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pov.get("align", align)
        p.space_after = Pt(pov.get("space", space))
        if lh:
            p.line_spacing = lh
        for part in runs:
            ov = {}
            if isinstance(part, tuple):
                part, ov = part
            r = p.add_run()
            r.text = part
            font(r, ov.get("name", name), ov.get("size", size),
                 ov.get("bold", bold), ov.get("color", color), ov.get("italic", False))
    return sh


def rect(slide, x, y, w, h, fill=None, line=BOX, lw=0.75, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.text_frame.text = ""
    return sh


def panel(slide, x, y, w, h, lines, size=14, fill=FILL, line=BOX, head=None,
          head_size=15, color=INK, name=JP, pad=0.14, anchor=MSO_ANCHOR.TOP, space=6):
    rect(slide, x, y, w, h, fill, line)
    dy = 0
    if head:
        rect(slide, x, y, w, 0.48, FILL2, line)
        text(slide, x, y, w, 0.48, head, size=head_size, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, pad=0.14)
        dy = 0.48
    if lines:
        text(slide, x, y + dy, w, h - dy, lines, size=size, color=color, name=name,
             pad=pad, anchor=anchor, space=space)


def code(slide, x, y, w, h, lines, size=13, head=None):
    rect(slide, x, y, w, h, CODEBG, RULE)
    dy = 0
    if head:
        text(slide, x, y + 0.06, w, 0.38, head, size=13, bold=True, color=MUTED, pad=0.14)
        dy = 0.42
    body = []
    for ln in lines:
        ov = {"space": 0, "name": MONO, "size": size}
        if isinstance(ln, tuple):
            ln, extra = ln
            ov.update(extra)
        body.append((ln, ov))
    text(slide, x, y + dy, w, h - dy, body, size=size, name=MONO, pad=0.14, space=0)


def arrow(slide, x, y, w, h, shape=MSO_SHAPE.LEFT_RIGHT_ARROW, label=None, size=15):
    sh = rect(slide, x, y, w, h, FILL2, BOX, shape=shape)
    if label:
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        font(r, JP, size, True, INK)
    return sh


def hline(slide, x, y, w, color=RULE, lw=1.0):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(lw)
    return ln


def photo(slide, x, y, w, h, path, caption=None):
    from PIL import Image
    iw, ih = Image.open(path).size
    inner_h = h - (0.32 if caption else 0.0)
    scale = min(w / iw, inner_h / ih)
    pw, ph = iw * scale, ih * scale
    slide.shapes.add_picture(path, Inches(x + (w - pw) / 2), Inches(y + (inner_h - ph) / 2),
                             Inches(pw), Inches(ph))
    rect(slide, x + (w - pw) / 2, y + (inner_h - ph) / 2, pw, ph, None, RULE)
    if caption:
        text(slide, x, y + inner_h, w, 0.32, caption, size=12, color=FAINT,
             align=PP_ALIGN.CENTER, pad=0.0)


def plain_table(slide, x, y, w, h, data, widths=None, size=13.5, head_size=13.5,
                row_h=0.45, aligns=None):
    gf = slide.shapes.add_table(len(data), len(data[0]),
                                Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    for e in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(e)
    sid = tblPr.makeelement(qn("a:tableStyleId"), {})
    sid.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"
    tblPr.append(sid)
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    if widths:
        tot = sum(widths)
        for i, cw in enumerate(widths):
            tbl.columns[i].width = Emu(int(Inches(w) * cw / tot))
    for r in range(len(data)):
        tbl.rows[r].height = Inches(row_h)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = FILL2 if r == 0 else WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            ov = {}
            if isinstance(val, tuple):
                val, ov = val
            p = tf.paragraphs[0]
            p.alignment = ov.get("align", (aligns[c] if aligns else PP_ALIGN.LEFT))
            rn = p.add_run()
            rn.text = val
            font(rn, ov.get("name", JP), ov.get("size", head_size if r == 0 else size),
                 ov.get("bold", r == 0), ov.get("color", INK))
    return tbl


TOTAL_PAGES = 15
_page = [0]


def slide(title=None, note=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        _page[0] += 1
        text(s, M, 0.28, CW - 2.0, 0.66, title, size=28, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, pad=0.0)
        text(s, W - M - 3.4, 0.34, 3.4, 0.54, "%d / %d" % (_page[0], TOTAL_PAGES),
             size=13, color=FAINT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, pad=0.0)
        hline(s, M, 1.06, CW)
        text(s, W - M - 1.2, H - 0.44, 1.2, 0.3, str(_page[0]), size=11, color=FAINT,
             align=PP_ALIGN.RIGHT, pad=0.0)
    if note:
        s.notes_slide.notes_text_frame.text = note
    return s


# ------------------------------------------------------------------ title
s = slide(note="つかみ:「イヤホンはコンピュータである」。BES2300YP は Cortex-M4F が 2 個載った"
               "れっきとした MCU。音を鳴らしていない時間、CPU は暇をしている。")
rect(s, M, 2.30, CW, 0.07, INK, None)
text(s, M, 2.58, CW, 1.30, "耳で並列計算してみた", size=58, bold=True, pad=0.0)
text(s, M, 3.98, CW, 0.56, "夏休みの自由研究 — ワイヤレスイヤホン 2 個を MPI クラスタにする",
     size=23, color=MUTED, pad=0.0)
rect(s, M, 4.72, CW, 0.04, RULE, None)
text(s, M, 5.05, 8.0, 1.2,
     [("[ 発表者名 ] / [ 所属 ]", {"size": 18, "space": 6}),
      ("github.com/TamichiRyuto/pine-buds-cluster", {"size": 16, "name": MONO,
                                                     "color": MUTED})], pad=0.0)
text(s, W - M - 4.0, 5.05, 4.0, 0.6, "LT 8 分", size=18, color=MUTED,
     align=PP_ALIGN.RIGHT, pad=0.0)

# ------------------------------------------------------------------ 1
s = slide("自己紹介 & サマリ",
          note="この 3 行と 6 個の数字だけ持ち帰ってもらえれば成功。"
               "「速くなった」とは一言も言っていないことに注意。")
panel(s, M, 1.26, 4.10, 2.66, [
    ("名前:   [                ]", {"space": 12}),
    ("所属:   [                ]", {"space": 12}),
    ("普段:   [                ]", {"space": 12}),
    ("GitHub: TamichiRyuto", {"space": 0}),
], size=16, head="発表者")
panel(s, 4.92, 1.26, CW - 4.37, 2.66, [
    ("1.  ワイヤレスイヤホン 2 個を、2 ノードの計算クラスタにした", {"space": 12}),
    ("2.  MPI を自作し、イヤホン同士の Bluetooth リンクに載せた", {"space": 12}),
    ("3.  標準 MPI API のベンチを、1 行も変えずに実機で PASS させた", {"space": 4}),
    ("     (OpenMP は API 互換のスタブ。ノード内 2 コア並列はまだ)",
     {"space": 0, "size": 13, "color": MUTED}),
], size=17, head="3 行サマリ")

tiles = [
    ("2 ノード", "右 = rank 0 / 左 = rank 1"),
    ("checksum = 32768", "float32 で厳密一致 → PASS"),
    ("max_payload 512 B", "ヘッダ 300 / バイナリ 672 / 実測 512"),
    ("1,788 checks", "実機に焼く前のホストテスト"),
    ("5 連タップ", "耳を叩くと両バッズで再実行"),
    ("80+ commits", "2026-08-18 → 09-01 の 2 週間"),
]
tw, gap = (CW - 2 * 0.22) / 3, 0.22
for i, (big, small) in enumerate(tiles):
    tx = M + (i % 3) * (tw + gap)
    ty = 4.22 + (i // 3) * 1.34
    rect(s, tx, ty, tw, 1.18, FILL, BOX)
    text(s, tx, ty + 0.10, tw, 0.56, big, size=22, bold=True, align=PP_ALIGN.CENTER, pad=0.05)
    text(s, tx, ty + 0.68, tw, 0.44, small, size=13, color=MUTED,
         align=PP_ALIGN.CENTER, pad=0.05)

# ------------------------------------------------------------------ 2
s = slide("並列計算のおさらい — MPI と OpenMP",
          note="HPC の定石「ノード間は MPI、ノード内は OpenMP」。"
               "今回はこの標準的な書き方をそのまま持ち込むのがテーマ。OpenMP 側は API だけ揃えた段階。")
NY, NW, NH = 1.24, 5.35, 2.72
for i, (nm, mem) in enumerate([("ノード A", "メモリ空間 A"), ("ノード B", "メモリ空間 B")]):
    bx = M + i * (CW - NW)
    rect(s, bx, NY, NW, NH, WHITE, BOX)
    text(s, bx, NY + 0.08, NW, 0.42, "%s  (%s)" % (nm, mem), size=17, bold=True,
         align=PP_ALIGN.CENTER, pad=0.05)
    for c in range(2):
        cx = bx + 0.40 + c * 2.45
        rect(s, cx, NY + 0.58, 2.10, 0.68, FILL2, BOX)
        text(s, cx, NY + 0.58, 2.10, 0.68, "コア %d" % c, size=16, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, pad=0.03)
    rect(s, bx + 0.40, NY + 1.36, 4.55, 0.46, FILL, BOX)
    text(s, bx + 0.40, NY + 1.36, 4.55, 0.46, "同じ配列を共有 (共有メモリ)", size=14,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, pad=0.03)
    text(s, bx + 0.25, NY + 1.90, 4.85, 0.62,
         "スレッドレベル並列 = OpenMP", size=14, color=MUTED,
         align=PP_ALIGN.CENTER, pad=0.03)
arrow(s, 6.06, NY + 0.80, 1.21, 0.68, MSO_SHAPE.LEFT_RIGHT_ARROW, "MPI", 17)
text(s, 5.90, NY + 1.56, 1.53, 0.8, "メモリが\n物理的に別", size=13, color=MUTED,
     align=PP_ALIGN.CENTER, pad=0.02, space=2)
text(s, M, 3.98, CW, 0.40,
     "クラスタレベル並列 = MPI :  明示的にメッセージを送り合う (MPI_Send / MPI_Recv / MPI_Barrier)",
     size=14.5, color=MUTED, align=PP_ALIGN.CENTER, pad=0.0)

plain_table(s, M, 4.44, 6.30, 1.70, [
    ["", "OpenMP", "MPI"],
    ["メモリ", "共有", "分散 (別空間)"],
    ["コスト", "安い (ns 〜 μs)", "高い (μs 〜 ms)"],
    ["今回の担当", "2 コア並列 (未達)", "イヤホン間 (左 ⇄ 右)"],
], widths=[1.5, 2.6, 2.2], size=14, head_size=14, row_h=0.42)

BX, BY, BW = 7.22, 4.44, 5.56
rect(s, BX, BY, BW, 2.14, WHITE, BOX)
text(s, BX, BY + 0.06, BW, 0.36, "GEMM (行列積) を 2 ランクに分ける", size=15, bold=True,
     align=PP_ALIGN.CENTER, pad=0.05)
text(s, BX + 0.14, BY + 0.44, BW - 0.28, 0.60, [
    [("N", {"name": MONO, "bold": True}), " = 行列の一辺。N=32 なら 32×32 の行列"],
    [("size", {"name": MONO, "bold": True}), " = MPI のランク数 = 参加したバッズの数"],
], size=13.5, pad=0.03, space=4)
rect(s, BX + 0.30, BY + 1.08, 2.45, 0.46, FILL2, BOX)
text(s, BX + 0.30, BY + 1.08, 2.45, 0.46, "rank 0 → 行 0 .. 15", size=14,
     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.03)
rect(s, BX + 2.85, BY + 1.08, 2.45, 0.46, FILL, BOX)
text(s, BX + 2.85, BY + 1.08, 2.45, 0.46, "rank 1 → 行 16 .. 31", size=14,
     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.03)
text(s, BX + 0.12, BY + 1.62, BW - 0.24, 0.44,
     "size=1 に縮退しても同じ PASS が出る → size を必ず見る",
     size=12.5, color=MUTED, align=PP_ALIGN.CENTER, pad=0.03)

# ------------------------------------------------------------------ 3
s = slide("今回のクラスタ構成",
          note="「USB を挿したまま」ではなく「耳に入れたまま」に近づいたことを強調する。"
               "rank は固定、IBRT のロールは起動ごとに入れ替わる。")
for i, (nm, rk, role) in enumerate([("右バッズ", "rank 0", "nv_role = MASTER"),
                                    ("左バッズ", "rank 1", "nv_role = SLAVE")]):
    bx = 1.05 + i * 6.60
    rect(s, bx, 1.24, 4.62, 1.76, WHITE, BOX)
    text(s, bx, 1.32, 4.62, 0.46, "%s   %s" % (nm, rk), size=20, bold=True,
         align=PP_ALIGN.CENTER, pad=0.05)
    text(s, bx, 1.82, 4.62, 1.10,
         "BES2300YP  /  %s\nCortex-M4F ×2  (単精度 FPU)" % role,
         size=14, color=MUTED, align=PP_ALIGN.CENTER, pad=0.05, space=4)
arrow(s, 5.88, 1.82, 1.56, 0.62, MSO_SHAPE.LEFT_RIGHT_ARROW, "IBRT", 16)
text(s, 5.72, 2.50, 1.88, 0.56, "左右間の\nTWS 制御チャネル",
     size=12.5, color=MUTED, align=PP_ALIGN.CENTER, pad=0.02, space=1)
text(s, M, 3.10, CW, 0.30,
     "※ rank は左右で固定。IBRT のロール (master / slave) は起動ごとに入れ替わる",
     size=12.5, color=MUTED, align=PP_ALIGN.CENTER, pad=0.0)
rect(s, 1.05, 3.46, 11.22, 0.56, FILL, BOX)
text(s, 1.05, 3.46, 11.22, 0.56,
     "Bluetooth SPP (自作ログチャネル)  =  Bluetooth 上のシリアルポート。UART とは独立",
     size=15, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.05)
arrow(s, 6.44, 4.10, 0.44, 0.36, MSO_SHAPE.DOWN_ARROW)
rect(s, 4.30, 4.56, 4.72, 0.98, WHITE, BOX)
text(s, 4.30, 4.62, 4.72, 0.38, "Windows PC   (COM6)", size=15, bold=True,
     align=PP_ALIGN.CENTER, pad=0.05)
text(s, 4.30, 5.00, 4.72, 0.44, "#23 GEMM-MPI ... PASS", size=14, name=MONO,
     color=MUTED, align=PP_ALIGN.CENTER, pad=0.05)
cols = [("ノード間 = MPI", "自作 MPI を IBRT に載せた"),
        ("ノード内 = OpenMP", "逐次スタブ。2 コア目は未使用"),
        ("観測 = SPP", "seq 番号で欠落・重複を検出")]
tw = (CW - 2 * 0.22) / 3
for i, (h, b) in enumerate(cols):
    tx = M + i * (tw + 0.22)
    panel(s, tx, 5.72, tw, 1.10, [(b, {"space": 0})], size=13.5, head=h, head_size=14.5,
          color=MUTED)

# ------------------------------------------------------------------ 4
s = slide("実験機材: PineBuds Pro",
          note="手元の個体のケースは Wiki 記載の CH342DS ではなく CH347 だった (ハードリビジョン差)。"
               "「完全オープンソース」とは言わない — IBRT/BT スタック中核はソースが無い。")
plain_table(s, M, 1.26, 7.35, 3.60, [
    ["項目", "値"],
    ["SoC", "Bestechnic BES2300YP (左右に 1 個ずつ、独立)"],
    ["CPU", "Dual-core Cortex-M4F @ 最大 300 MHz / 単精度 FPU のみ"],
    ["メモリ", "SRAM 992 KB / Flash 4 MB"],
    ["ファーム", "OpenPineBuds — アプリ層は OSS (BT 中核は .a)"],
    ["デバッグ", "充電ケースが USB → デュアル UART (2 Mbaud)"],
    ["位置づけ", "実売 1 万円前後、普通に音楽が聴ける TWS イヤホン"],
], widths=[1.35, 6.0], size=14.5, head_size=14.5, row_h=0.50)
if PHOTO:
    photo(s, 8.20, 1.26, 4.58, 2.20, PHOTO,
          "PineBuds Pro 本体と充電ケース (ケースが USB → デュアル UART プログラマ)")
else:
    rect(s, 8.20, 1.26, 4.58, 2.20, FILL, RULE)
    text(s, 8.20, 1.26, 4.58, 2.20,
         "[  PineBuds Pro の写真をここに  ]\nslides/assets/pinebuds.jpg を置いて\n"
         "scripts/make_lt_pptx.py を再実行",
         size=14, color=FAINT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         pad=0.1, space=4)
panel(s, 8.20, 3.62, 4.58, 1.90, [
    ("ケースに USB-C を挿すと", {"space": 6}),
    ("/dev/ttyACM0 (右) と ttyACM1 (左)", {"space": 6, "name": MONO, "size": 13}),
    ("が生えて、自作ファームが焼ける。", {"space": 10}),
    ("イヤホンなのに開発ボード。", {"space": 0, "bold": True}),
], size=14.5, head="つまり")
panel(s, M, 5.20, 7.35, 1.55, [
    ("・ 単コア GEMM が 5 〜 10 ms で回る計算資源が、耳の中に 2 個ある", {"space": 6}),
    ("・ RAM は約 330 KB 空き (2026-08-30 ビルドの .map 集計)", {"space": 0}),
], size=14.5, head="なぜ計算機として使えるのか")

# ------------------------------------------------------------------ 5
s = slide("なぜ「普通のワイヤレスイヤホン」ではダメだったのか",
          note="「Linux が動くイヤホン」ではなく「自分のコードを載せられるイヤホン」を探した、という話。"
               "文鎮化リスクは共感を得やすい。最初にやったのは工場ファームの全量バックアップ。")
rect(s, M, 1.22, CW, 0.60, FILL, BOX)
text(s, M, 1.22, CW, 0.60,
     "当初は手持ちの TWS でやりたかった。 →  自作コードを載せる手段が見つからなかった。",
     size=17, bold=True, anchor=MSO_ANCHOR.MIDDLE, pad=0.14)
plain_table(s, M, 1.96, CW, 3.20, [
    ["自作コードをイヤホンで動かすのに必要なもの", "調べた範囲の市販 TWS", "PineBuds Pro"],
    ["① ファームウェアのソース / SDK",
     ("✗  非公開", {"color": NG, "bold": True}),
     ("✓  OSS (OpenPineBuds)", {"color": OK, "bold": True})],
    ["② 署名されていない FW を焼ける経路",
     ("✗  手段が非公開", {"color": NG, "bold": True}),
     ("✓  bestool で焼ける", {"color": OK, "bold": True})],
    ["③ デバッグ UART が外に出ている",
     ("✗  出ていない", {"color": NG, "bold": True}),
     ("✓  充電ケースが兼ねる", {"color": OK, "bold": True})],
    ["④ 左右間リンクを叩ける API",
     ("✗  非公開の独自 TWS", {"color": NG, "bold": True}),
     ("✓  ヘッダで見える", {"color": OK, "bold": True})],
    ["⑤ 文鎮化したときの復旧手段",
     ("✗  見つからなかった", {"color": NG, "bold": True}),
     ("✓  工場 FW + 純正ライタ", {"color": OK, "bold": True})],
], widths=[5.6, 3.1, 3.5], size=14.5, head_size=14, row_h=0.50)
text(s, M, 5.30, CW, 0.80, [
    ("・ ① 〜 ③ のどれか 1 つでも欠けると、コードを書いても載せられない = 実験が成立しない", {"space": 7}),
    ("・ これは手元の機種と調べられた範囲の話で、すべての TWS がそうだと主張するものではない",
     {"space": 0, "size": 13.5}),
], size=14.5, color=MUTED)
rect(s, M, 6.20, CW, 0.66, FILL2, BOX)
text(s, M, 6.20, CW, 0.66,
     "→  この 5 つが揃う TWS が、探した範囲でここだけだったので購入",
     size=17, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.14)

# ------------------------------------------------------------------ 6
s = slide("やったこと — 全体像",
          note="どこまでが標準 API で、どこからが自作かを一目で見せる。"
               "bench/ は本物の OpenMPI でもそのままビルドできる。")
layers = [
    ("bench/gemm_mpi_omp.cpp", "標準 MPI + OpenMP API で書いた GEMM。実機用に 1 行も書き換えない ← これが要件", FILL2),
    ("adapters/mpi/  +  adapters/omp/", "自作 MPI サブセット (Send / Recv / Barrier / Allreduce / Isend / Wait) と OpenMP スタブ", FILL),
    ("transport seam (関数ポインタ 1 枚)", "ホスト = pthread ループバック (TDD 用) / 実機 = 断片化 + クレジット制御", FILL),
    ("firmware/pinebuds_compute/", "IBRT glue / SPP ログチャネル / 5 連タップ起動の調停", FILL),
    ("OpenPineBuds SDK", "IBRT (TWS 制御) · BESAUD · RTX (CMSIS-RTOS v1)", FILL2),
]
for i, (h, b, f) in enumerate(layers):
    ly = 1.22 + i * 0.92
    rect(s, M, ly, CW, 0.84, f, BOX)
    text(s, M + 0.14, ly + 0.06, CW - 0.28, 0.38, h, size=17, bold=True, pad=0.02)
    text(s, M + 0.14, ly + 0.42, CW - 0.28, 0.38, b, size=13.5, color=MUTED, pad=0.02)
panel(s, M, 5.90, CW, 1.00, [
    ("・ gnu++98 / 例外・RTTI 禁止   ・ ヒープ禁止・STL 禁止・静的バッファのみ   "
     "・ double 禁止 (単精度 FPU)", {"space": 0}),
], size=14.5, head="SDK に合わせるしかない制約")

# ------------------------------------------------------------------ 7
LANE_FILL = [RGBColor(0xE9, 0xF1, 0xF8), RGBColor(0xE4, 0xF2, 0xF3),
             RGBColor(0xEE, 0xEB, 0xF7), RGBColor(0xF8, 0xF2, 0xE6)]
LANE_HEAD = [RGBColor(0xD6, 0xE6, 0xF3), RGBColor(0xCF, 0xE8, 0xEA),
             RGBColor(0xE1, 0xDC, 0xF1), RGBColor(0xF2, 0xE6, 0xCF)]
s = slide("開発環境 — 1 回の実験に必要なもの",
          note="4 つの世界をまたがないと 1 回も試せない、という絵。"
               "「WSL の中の Docker の中の GCC で焼いて、Windows の Python で受ける」で笑いが取れる。")
lanes = [
    ("🐧", "WSL2  Ubuntu", ["コードを書くのは全部ここ",
                            "make test → 1,788 checks",
                            "install-into-sdk.sh で SDK へ"]),
    ("🐳", "Docker  SDK", ["./start_dev.sh でコンテナ",
                           "./build.sh → open_source.bin",
                           "backup.sh で工場 FW を退避"]),
    ("🪟", "Windows 11", ["usbipd でケースを WSL へ",
                          "→ /dev/ttyACM0 と ttyACM1",
                          "spp_tail.py COM6 で SPP 受信"]),
    ("🎧", "PineBuds Pro", ["bestool write-image を左右に",
                            "出して 3 秒待って戻す = リブート",
                            "picocom -b 2000000 で観測"]),
]
lw = (CW - 3 * 0.20) / 4
for i, (icon, name, items) in enumerate(lanes):
    lx = M + i * (lw + 0.20)
    rect(s, lx, 1.24, lw, 2.14, LANE_FILL[i], BOX)
    rect(s, lx, 1.24, lw, 0.60, LANE_HEAD[i], BOX)
    text(s, lx + 0.10, 1.30, 0.52, 0.48, icon, size=22, pad=0.0)
    text(s, lx + 0.66, 1.24, lw - 0.76, 0.60, name, size=15, bold=True,
         anchor=MSO_ANCHOR.MIDDLE, pad=0.0)

    def _jp(t):
        return any("぀" <= c <= "ヿ" or "一" <= c <= "鿿" for c in t)
    body = [(t, {"space": 9, "size": 12, "name": JP if _jp(t) else MONO}) for t in items]
    text(s, lx + 0.08, 1.92, lw - 0.16, 1.40, body, size=12, pad=0.05)

CC_FILL, CC_HEAD = RGBColor(0xF3, 0xEF, 0xF9), RGBColor(0xE3, 0xDC, 0xF2)
rect(s, M, 3.54, CW, 1.10, CC_FILL, BOX)
rect(s, M, 3.54, CW, 0.48, CC_HEAD, BOX)
text(s, M + 0.10, 3.56, 0.50, 0.44, "🤖", size=19, pad=0.0)
text(s, M + 0.64, 3.54, 6.4, 0.48, "Claude Code — 4 つの世界を横断して回す", size=15,
     bold=True, anchor=MSO_ANCHOR.MIDDLE, pad=0.0)
text(s, W - M - 5.4, 3.54, 5.4, 0.48, "103 コミット中 102 が Co-Authored-By: Claude",
     size=12, color=MUTED, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, pad=0.12)
cc = [("設計・SDK 調査", "プリビルト .a の逆アセンブルと nm"),
      ("TDD の運用", "Red は本体、Green はサブエージェント"),
      ("実機ランの監視", "UART / SPP を監視して docs に反映")]
ccw = (CW - 0.40 - 2 * 0.20) / 3
for i, (h, b) in enumerate(cc):
    cx = M + 0.20 + i * (ccw + 0.20)
    text(s, cx, 4.06, ccw, 0.28, "・ " + h, size=13, bold=True, pad=0.0)
    text(s, cx + 0.20, 4.32, ccw - 0.20, 0.28, b, size=12, color=MUTED, pad=0.0)

rect(s, M, 4.80, CW, 1.62, FILL, BOX)
rect(s, M, 4.80, CW, 0.48, FILL2, BOX)
text(s, M, 4.80, CW, 0.48, "1 サイクル — 4 つの世界を 1 周しないと 1 回も試せない", size=15,
     bold=True, anchor=MSO_ANCHOR.MIDDLE, pad=0.14)
chips = ["① 書く", "② ホストで緑", "③ SDK に配置", "④ ビルド", "⑤ 焼く", "⑥ 耳で観る"]
cw_ = (CW - 0.60 - 5 * 0.28) / 6
for i, c in enumerate(chips):
    cx = M + 0.30 + i * (cw_ + 0.28)
    rect(s, cx, 5.38, cw_, 0.60, WHITE, BOX)
    text(s, cx, 5.38, cw_, 0.60, c, size=13.5, bold=True, anchor=MSO_ANCHOR.MIDDLE,
         align=PP_ALIGN.CENTER, pad=0.02)
    if i < 5:
        arrow(s, cx + cw_ + 0.01, 5.53, 0.26, 0.30, MSO_SHAPE.RIGHT_ARROW)
text(s, M, 6.04, CW, 0.34,
     "② が緑にならないうちは ④ に進まない — 焼いてから気づくと左右 2 個ぶんやり直し",
     size=13, color=MUTED, align=PP_ALIGN.CENTER, pad=0.0)
text(s, M, 6.56, CW, 0.40,
     "WSL の中の Docker の中の GCC でビルドして、Windows の Python で結果を受ける",
     size=15, bold=True, align=PP_ALIGN.CENTER, pad=0.0)

# ------------------------------------------------------------------ 8
s = slide("実装の山場 ① — SDK にソースが無い",
          note="「ヘッダを信じるな、リンクされるバイナリを信じろ」。"
               "300 / 672 / 512 の 3 つの数字を並べるのがこのスライドの山。時間を使うところ。")
rect(s, M, 1.22, CW, 0.66, FILL, BOX)
text(s, M, 1.22, CW, 0.66,
     "問題: 左右リンクの本体は inc/ と lib/ しか無い = プリビルトの .a しか無い → 仕様が読めない",
     size=16.5, bold=True, anchor=MSO_ANCHOR.MIDDLE, pad=0.14)
rect(s, M, 2.00, CW, 0.60, FILL2, BOX)
text(s, M, 2.00, CW, 0.60,
     "やったこと: 静的ライブラリを逆アセンブルして、仕様を実物から確定した",
     size=16.5, bold=True, anchor=MSO_ANCHOR.MIDDLE, pad=0.14)
plain_table(s, M, 2.74, CW, 1.60, [
    ["調べたこと", "ヘッダの記述", "逆アセンブルした実体"],
    ["送信ペイロード上限", "300", ("実行時 assert は 672", {"bold": True})],
    ["RX ハンドラのスレッド", "記述なし", ("BT スタック本体のループ (止めると全部止まる)", {"bold": True})],
], widths=[3.2, 3.2, 5.8], size=15, head_size=14.5, row_h=0.52)
nums = [("300", "ヘッダの定数"), ("672", "バイナリの assert"), ("512", "実測値")]
nw = (CW - 2 * 0.30) / 3
for i, (n, cap) in enumerate(nums):
    nx = M + i * (nw + 0.30)
    rect(s, nx, 4.56, nw, 1.10, FILL if i < 2 else FILL2, BOX)
    text(s, nx, 4.62, nw, 0.62, n, size=34, bold=True, align=PP_ALIGN.CENTER, pad=0.02)
    text(s, nx, 5.24, nw, 0.34, cap, size=13.5, color=MUTED, align=PP_ALIGN.CENTER, pad=0.02)
text(s, M, 5.82, CW, 0.44,
     "起動時に 4 → 64 → … → 668 B を掃引して上限を実測。256 B 超はフレームに分割し、"
     "in-flight を 536 B に抑える",
     size=14, color=MUTED, align=PP_ALIGN.CENTER, pad=0.0)
rect(s, M, 6.34, CW, 0.56, FILL2, BOX)
text(s, M, 6.34, CW, 0.56, "ヘッダを信じるな。リンクされるバイナリを信じろ。",
     size=17, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.1)

# ------------------------------------------------------------------ 9
s = slide("実装の山場 ② — 実機は 1 発では動かない",
          note="「組み込みは電源とペアリングで死ぬ」は共感ポイント。"
               "MMFAR の値が構造体オフセットと一致した話は必ず言う。")
plain_table(s, M, 1.22, 7.70, 4.10, [
    ["#", "症状", "原因"],
    ["1", "finalize の 2 秒後に MemFault で無限リブート", "RTX でスレッドが自分自身を終了させると壊れる"],
    ["2", "充電起動すると BT スタックが上がらない", "充電イベントのリセット設定の副作用"],
    ["3", "起動 15 〜 60 秒で勝手に電源が落ちる", "満充電シャットダウン"],
    ["4", "PC と再ペアリングできない", "通常ビルドは inquiry scan を出さない"],
    ["5", "SPP ログの先頭 512 B だけ二重に届く", "1 回の接続で CONNECTED が 2 回上がる"],
], widths=[0.5, 3.7, 4.4], size=14, head_size=14, row_h=0.62)
code(s, 8.50, 1.22, 4.28, 2.60, [
    "### EXCEPTION ###",
    ("MMFAR=00000034", {"bold": True}),
    "FaultCause: (Data access violation)",
    ("Current Task    : 0", {"bold": True}),
    "New Running Task: 255",
], size=13.5, head="Run 2 : 毎ブート・17.5 秒周期")
panel(s, 8.50, 3.98, 4.28, 1.60, [
    ("0x34 = RTOS のスレッド構造体の", {"space": 5}),
    ("フィールドのオフセットと一致。", {"space": 9}),
    ("終了済みスレッドを触っていた。", {"space": 0, "bold": True}),
], size=14, head="MMFAR が犯人を名指しした")
rect(s, M, 5.68, 7.70, 0.66, FILL2, BOX)
text(s, M, 5.68, 7.70, 0.66, "計算に起因する地雷はゼロ。全部が電源・接続・RTOS。",
     size=16, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.1)
text(s, 8.50, 5.76, 4.28, 0.50,
     "純ロジックは全部ホストのテストを先に書いた", size=13.5, color=MUTED, pad=0.0)

# ------------------------------------------------------------------ 10
s = slide("実装の山場 ③ — 観測をワイヤレスに / 耳で操作する",
          note="デモ的に一番「おっ」となるところ。ケースの中でもタッチは反応する。")
rect(s, M, 1.22, CW, 0.54, FILL, BOX)
text(s, M, 1.22, CW, 0.54, "USB ケーブルで繋がったイヤホンは、もはやイヤホンではない",
     size=17, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.1)
panel(s, M, 1.94, 6.05, 2.30, [
    ("・ ログをリングバッファに積み、Bluetooth の", {"space": 4}),
    ("   シリアルポート (SPP) で PC に流す", {"space": 9, "color": MUTED}),
    ("・ seq 番号で欠落・重複を機械検出", {"space": 9}),
    ("・ UART とは独立の経路 = 原理的に USB は不要", {"space": 0, "bold": True}),
], size=14.5, head="① Bluetooth SPP に自作ログチャネル")
panel(s, 6.78, 1.94, 6.00, 2.30, [
    ("・ タッチパッドは電源キー入力。ジェスチャは", {"space": 4}),
    ("   SDK がソフトで判定している", {"space": 9, "color": MUTED}),
    ("・ 単タップ 〜 4 連・長押しは既存機能で満杯", {"space": 9}),
    ("・ 空いていた 5 連タップを奪った", {"space": 0, "bold": True}),
], size=14.5, head="② 5 連タップで GEMM を再実行")
FY = 4.52
steps = [("① 5 連タップ", "叩かれた側だけが\nキーイベントを受け取る"),
         ("② 相手に START", "左右間の制御チャネルに\n5 バイトを 1 本"),
         ("③ 相手も起きる", "待機セマフォを解放して\n計算に入る"),
         ("④ 両バッズで GEMM", "rank 0 が checksum を\n検証して PASS")]
sw = 2.72
gap = (CW - 4 * sw) / 3
for i, (h, b) in enumerate(steps):
    sx = M + i * (sw + gap)
    rect(s, sx, FY, sw, 1.56, FILL if i % 2 == 0 else WHITE, BOX)
    text(s, sx, FY + 0.10, sw, 0.40, h, size=15, bold=True, align=PP_ALIGN.CENTER, pad=0.04)
    text(s, sx, FY + 0.54, sw, 0.94, b, size=12.5, color=MUTED,
         align=PP_ALIGN.CENTER, pad=0.04, space=2)
    if i < 3:
        arrow(s, sx + sw + 0.06, FY + 0.60, gap - 0.12, 0.38, MSO_SHAPE.RIGHT_ARROW)
text(s, M, 6.28, CW, 0.42,
     "実行中の連打は捨てる。相手の START が落ちたら、タップした側だけが縮退して完走する。",
     size=13.5, color=MUTED, align=PP_ALIGN.CENTER, pad=0.0)

# ------------------------------------------------------------------ 11
s = slide("結果 ① — 実機ログ (両バッズ / 全条件 PASS)",
          note="「32768 = 32³。全要素 1 の行列積は整数値しか通らないので float32 で厳密に一致する。"
               "誤差 1 ulp も許さない判定」は必ず言う。")
code(s, M, 1.22, 7.55, 4.20, [
    ("== 右バッズ (rank 0) ==", {"bold": True}),
    "[mpi] init rank=0 size=2 link_wait=200 ms",
    "[mpi] peer ok rank=0 peer=1",
    "[mpi] barrier ok",
    ("GEMM-MPI N=32 rank=0 size=2", {"bold": True}),
    ("  checksum=32768.000000", {"bold": True}),
    ("  expect  =32768.000000  PASS", {"bold": True}),
    ("GEMM-MPI elapsed=13 ms tx=2 rx=4 err=0", {"bold": True}),
    "",
    ("== 左バッズ (rank 1) ==", {"bold": True}),
    "[mpi] init rank=1 size=2 link_wait=500 ms",
    "[mpi] send ok / frames tx=2 rx=2 err=0",
    "[mpi] finalize done rank=1",
], size=13.5, head="Run 4 — ケース内・充電起動・両バッズ同時")
checks = [
    "両側 size=2 — 本当に 2 ノードで走った",
    "checksum=32768.000000 の厳密一致",
    "rank 0 / rank 1 が左右で 1 つずつ",
    "EXCEPTION 0 回・送信失敗 0 回",
]
rect(s, 8.32, 1.22, 4.46, 2.80, FILL, BOX)
rect(s, 8.32, 1.22, 4.46, 0.48, FILL2, BOX)
text(s, 8.32, 1.22, 4.46, 0.48, "合否条件 — すべて充足", size=15.5, bold=True,
     anchor=MSO_ANCHOR.MIDDLE, pad=0.14)
for i, c in enumerate(checks):
    cy = 1.82 + i * 0.54
    text(s, 8.46, cy, 0.34, 0.44, "✓", size=18, bold=True, color=OK, pad=0.0)
    text(s, 8.84, cy - 0.01, 3.86, 0.50, c, size=14, pad=0.0)
panel(s, 8.32, 4.20, 4.46, 2.10, [
    ("32768 = 32³", {"space": 8, "bold": True, "size": 20, "name": MONO}),
    ("全要素 1 の行列積は整数値しか通らない", {"space": 5, "color": MUTED}),
    ("→ float32 でも丸め誤差がゼロ", {"space": 5, "color": MUTED}),
    ("→ 厳密一致以外を許さない判定にできる", {"space": 0, "bold": True}),
], size=14, head="なぜ checksum で判定できるのか")

# ------------------------------------------------------------------ 12
s = slide("結果 ② — ワイヤレス受信と 5 連タップ再実行",
          note="「耳を 5 回叩くと両バッズで行列積が走り、実行ログが Bluetooth で PC に流れてくる」で締める。"
               "PASS 行そのものが飛んだとは言わない。")
code(s, M, 1.22, 6.55, 2.70, [
    "#2  GEMM float N=32  checksum=32768.000000",
    "                     expect=32768.000000  PASS",
    "#7  [mpi] peer ok rank=1 peer=0",
    ("#19 [mpi-t1] max_payload=512", {"bold": True}),
    "#20 [mpi] barrier ok",
    "#23 [mpi] finalize done rank=1",
], size=13.5, head="Run 12 : Bluetooth で PC が受けたログ (欠落・重複ゼロ)")
plain_table(s, 7.30, 1.22, 5.48, 2.90, [
    ["run", "トリガ", "結果", "elapsed"],
    ["起動時", "自動", "PASS  size=2", "12 ms"],
    ["#1", "右を 5 連タップ", "PASS  size=2", "128 ms"],
    ["#2", "左を 5 連タップ", "PASS  size=2", "5 ms"],
    ["#3", "右を 5 連タップ", "PASS  size=2", "115 ms"],
    ["#4", "左を 5 連タップ", "PASS  size=2", "83 ms"],
], widths=[0.9, 2.0, 1.7, 1.1], size=14, head_size=14, row_h=0.46)
panel(s, M, 4.10, 6.55, 1.80, [
    ("タップした側が先に走り、相手は START の伝搬を待って", {"space": 5}),
    ("合流する。だからタップした側の elapsed が伸びる。", {"space": 0, "color": MUTED}),
], size=14, head="elapsed のばらつきの読み方")
panel(s, 7.30, 4.10, 5.48, 1.80, [
    ("このランで PC に出たのは左 (rank 1) のログだけ。", {"space": 5}),
    ("slave 側の送信はスタック上は成功扱いなのに電波に", {"space": 5, "color": MUTED}),
    ("乗らず、rank 0 の行は UART にしか出ない (未解決)。", {"space": 0, "color": MUTED}),
], size=14, head="正直な注記")
rect(s, M, 6.12, CW, 0.62, FILL2, BOX)
text(s, M, 6.12, CW, 0.62,
     "耳を 5 回叩くと両バッズで行列積が走り、実行ログが Bluetooth 経由で PC に流れてくる",
     size=17, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.1)

# ------------------------------------------------------------------ 13
s = slide("結果 ③ — 数字で見る",
          note="「速くなりませんでした」で終わらせず、なぜ遅いかが定量的に説明できることを成果にする。"
               "ここが最大の山場。表を指す前に口頭でオチを言う。")
plain_table(s, M, 1.22, 7.55, 1.95, [
    ["指標", "実測値", "コメント"],
    ["GEMM N=32  単コア", "5 〜 10 ms", "基準"],
    ["GEMM N=32  2 ノード", ("12 〜 13 ms", {"bold": True}), ("通信込み。遅くなった", {"bold": True})],
    ["アイドルからの往復", "最大 407 ms", "ランごとに大きく変動"],
], widths=[2.8, 2.2, 3.0], size=14.5, head_size=14, row_h=0.48)
rect(s, M, 3.42, 7.55, 1.14, FILL2, BOX)
text(s, M, 3.54, 7.55, 0.52, "単コア 5〜10 ms   →   2 ノード 12〜13 ms", size=21,
     bold=True, align=PP_ALIGN.CENTER, pad=0.05)
text(s, M, 4.06, 7.55, 0.36, "分散したぶんだけ遅くなった", size=14,
     align=PP_ALIGN.CENTER, pad=0.05)
panel(s, M, 4.74, 7.55, 1.58, [
    ("407 ms は 100 回連続プローブ中の値で、GEMM 実行中の", {"space": 5}),
    ("レイテンシではない。省電力状態から叩き起こすとここまで伸びる。",
     {"space": 0, "color": MUTED}),
], size=14, head="407 ms の読み方")
rect(s, 8.32, 1.22, 4.46, 5.10, FILL, BOX)
rect(s, 8.32, 1.22, 4.46, 0.48, FILL2, BOX)
text(s, 8.32, 1.22, 4.46, 0.48, "なぜ遅いのか = この実験の本題", size=15.5, bold=True,
     anchor=MSO_ANCHOR.MIDDLE, pad=0.14)
rows13 = [("① 分割で浮く計算", "65,000 flop を半分に\n→  数 ms"),
          ("② 往復に要る時間", "リンクが生きていても 数 ms\n→  ① とほぼ相殺"),
          ("③ アイドルからの往復", "最大 407 ms\n=  単コア GEMM 40 〜 80 回ぶん")]
for i, (h, b) in enumerate(rows13):
    ry = 1.86 + i * 1.10
    rect(s, 8.52, ry, 4.06, 1.00, WHITE, BOX)
    text(s, 8.52, ry + 0.06, 4.06, 0.34, h, size=14, bold=True, align=PP_ALIGN.CENTER, pad=0.04)
    text(s, 8.52, ry + 0.40, 4.06, 0.58, b, size=12.5, color=MUTED, align=PP_ALIGN.CENTER,
         pad=0.04, space=2)
text(s, 8.44, 5.28, 4.22, 0.92,
     "「計算量 / 通信量」が小さい問題を\n分散してはいけない", size=17, bold=True,
     align=PP_ALIGN.CENTER, pad=0.06, space=4)
rect(s, M, 6.40, CW, 0.58, FILL2, BOX)
text(s, M, 6.40, CW, 0.58,
     "速くはならなかった。が、HPC の教科書どおりの結論を、耳の中で定量的に再現できた。",
     size=16.5, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.1)

# ------------------------------------------------------------------ 14
s = slide("わかったこと / できていないこと",
          note="できていないことを先に自分で言い切ると、質疑が建設的になる。")
items_ok = [
    ("標準 API のまま実機に載る", "本物の OpenMPI でも自作アダプタでも無改変でビルドでき、結果も一致した"),
    ("ソースが無くても仕様は確定できる", "逆アセンブルで、ヘッダと食い違う実挙動を根拠付きで確定できた"),
    ("組み込みの敵は計算ではない", "踏んだ地雷は全部が電源・接続・RTOS・SDK の仕様だった"),
    ("ホストで TDD できる形に切ると強い", "実機で 1 回試す前に 1,788 checks。実機バグもテストに落として直した"),
]
items_ng = [
    ("ノード内 2 コア並列は未達", "OpenMP は逐次スタブのまま。2 個目のコアは音声処理と取り合いになる"),
    ("slave 側のログが電波に乗らない", "片方が master のときしか PC に届かない"),
    ("音楽再生との同時実行は未検証", "「使いながら裏で計算」はまだ道半ば"),
    ("N を大きくした時の挙動は未計測", "RAM の上限は N≈128。そこで逆転するかは未検証"),
]
for col, (head, items, mark, mc) in enumerate([
        ("わかったこと", items_ok, "✓", OK), ("できていないこと (正直に)", items_ng, "✗", NG)]):
    cx = M + col * (CW / 2 + 0.12)
    cw = CW / 2 - 0.12
    rect(s, cx, 1.22, cw, 5.44, FILL if col == 0 else WHITE, BOX)
    rect(s, cx, 1.22, cw, 0.52, FILL2, BOX)
    text(s, cx, 1.22, cw, 0.52, head, size=17, bold=True, anchor=MSO_ANCHOR.MIDDLE, pad=0.16)
    for i, (h, b) in enumerate(items):
        iy = 1.92 + i * 1.20
        text(s, cx + 0.16, iy, 0.36, 0.40, mark, size=17, bold=True, color=mc, pad=0.0)
        text(s, cx + 0.54, iy - 0.02, cw - 0.72, 0.40, h, size=14.5, bold=True, pad=0.0)
        text(s, cx + 0.54, iy + 0.38, cw - 0.72, 0.76, b, size=12.5, color=MUTED, pad=0.0)

# ------------------------------------------------------------------ 15
s = slide("まとめ",
          note="締め:「みなさんの耳にも、暇をしている Cortex-M4F が 2 個入っています」。")
rect(s, M, 1.22, CW, 0.92, FILL2, BOX)
text(s, M, 1.22, CW, 0.92, "1 万円のワイヤレスイヤホン 2 個  =  2 ノードの MPI クラスタ",
     size=27, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.1)
res = [(OK, "✓", "標準 MPI API のベンチが、無改変で耳の中で PASS した"),
       (OK, "✓", "checksum 厳密一致 / 両バッズ size=2 / EXCEPTION 0 回"),
       (OK, "✓", "5 連タップで再実行、結果は Bluetooth で PC にワイヤレス配信"),
       (NG, "✗", "速くはならなかった — が、理由は定量的に説明できる")]
for i, (c, m, t) in enumerate(res):
    ry = 2.44 + i * 0.58
    text(s, M + 0.20, ry, 0.38, 0.46, m, size=19, bold=True, color=c, pad=0.0)
    text(s, M + 0.64, ry + 0.02, CW - 0.9, 0.46, t, size=16.5, pad=0.0)
nexts = [("① ノード内 2 コア並列", "2 個目のコアを音声処理と\nどう分け合うかが論点"),
         ("② 通信に耐える問題を選ぶ", "境界だけを交換する問題なら\n通信量が小さい"),
         ("③ 音楽再生との同時実行", "本当に「使いながら裏で計算」\nさせる")]
tw = (CW - 2 * 0.22) / 3
for i, (h, b) in enumerate(nexts):
    tx = M + i * (tw + 0.22)
    panel(s, tx, 4.86, tw, 1.24, [(b, {"space": 0})], size=13, head=h, head_size=14,
          color=MUTED)
rect(s, M, 6.22, CW, 0.60, FILL2, BOX)
text(s, M, 6.22, CW, 0.60, "みなさんの耳にも、暇をしている Cortex-M4F が 2 個入っています",
     size=19, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, pad=0.1)
text(s, M, 6.92, CW, 0.44, "github.com/TamichiRyuto/pine-buds-cluster",
     size=17, bold=True, name=MONO, align=PP_ALIGN.CENTER, pad=0.0)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("wrote %s (%d slides)" % (OUT, len(prs.slides._sldIdLst)))
